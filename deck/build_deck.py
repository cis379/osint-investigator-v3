"""Build osint-brief.pptx from the briefing content — a Google-Slides-importable deck.

Build-only (python-pptx is NOT a system dependency; not in requirements.txt).
Prereq images: run `python deck/render_graphs.py` first (writes deck/img/*.png).
Run:  python deck/build_deck.py   ->   deck/osint-brief.pptx
Then: upload the .pptx to Google Drive and open with Google Slides (auto-converts).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

HERE = os.path.dirname(__file__)
IMG = os.path.join(HERE, "img")

# ---- palette (dark intelligence-briefing theme) ---------------------------
BG     = RGBColor(0x0D, 0x1B, 0x2A)
PANEL  = RGBColor(0x16, 0x32, 0x4F)
PANEL2 = RGBColor(0x10, 0x24, 0x3A)
BLUE   = RGBColor(0x58, 0xA6, 0xFF)
GREEN  = RGBColor(0x3F, 0xB9, 0x50)
AMBER  = RGBColor(0xD2, 0x99, 0x22)
PURPLE = RGBColor(0xBC, 0x8C, 0xFF)
PINK   = RGBColor(0xF7, 0x78, 0xBA)
TEXT   = RGBColor(0xE6, 0xED, 0xF3)
MUTE   = RGBColor(0x9F, 0xB0, 0xC3)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, dash=False):
    b = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        b.fill.background()
    else:
        b.fill.solid(); b.fill.fore_color.rgb = fill
    if line is None:
        b.line.fill.background()
    else:
        b.line.color.rgb = line; b.line.width = Pt(line_w)
        if dash:
            from pptx.oxml.ns import qn
            ln = b.line._get_or_add_ln(); d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'}); ln.append(d)
    b.shadow.inherit = False
    return b


def text(s, x, y, w, h, runs, size=16, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, space=4):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [[(runs, size, color, bold)]]
    elif runs and isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.space_before = Pt(0)
        if isinstance(para, tuple):
            para = [para]
        for (t, sz, col, bd) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bd; r.font.name = "Calibri"
    return tb


def bullets(s, x, y, w, h, items, size=17, color=TEXT, gap=7):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(0)
        if isinstance(it, str):
            it = [("•  " + it, size, color, False)]
        for (t, sz, col, bd) in it:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bd; r.font.name = "Calibri"
    return tb


def header(s, kicker, title):
    box(s, 0, 0, 13.333, 0.12, fill=BLUE)
    text(s, 0.55, 0.32, 12.2, 0.4, kicker.upper(), size=13, color=BLUE, bold=True)
    text(s, 0.55, 0.66, 12.2, 0.9, title, size=27, color=WHITE, bold=True)


def footer(s, name):
    text(s, 0.55, 7.08, 8, 0.32, "OSINT Investigator  ·  external-OSINT only  ·  runs locally", size=10, color=MUTE)
    text(s, 9.1, 7.08, 3.65, 0.32, name, size=10, color=MUTE, align=PP_ALIGN.RIGHT)


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt


def centered(b, label, size=12.5, color=TEXT, bold=True):
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(4); tf.margin_right = Pt(4); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, line in enumerate(label if isinstance(label, list) else [label]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(0); p.space_before = Pt(0)
        parts = line if isinstance(line, list) else [line]
        for (t, sz, col, bd) in parts:
            r = p.add_run(); r.text = t; r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bd; r.font.name = "Calibri"
    return b


def connect(s, x1, y1, x2, y2, color=BLUE, w=2.0):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(w)
    return c


def arrow_r(s, x, y, w=0.4, color=BLUE):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.26))
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background(); a.shadow.inherit = False
    return a


def arrow_d(s, x, y, h=0.35, color=BLUE):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(0.26), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background(); a.shadow.inherit = False
    return a


def pic(s, path, x, y, w=None, h=None):
    kw = {}
    if w: kw["width"] = Inches(w)
    if h: kw["height"] = Inches(h)
    return s.shapes.add_picture(path, Inches(x), Inches(y), **kw)


# ============================================================ SLIDE 1 — TITLE
s = slide()
box(s, 0, 0, 13.333, 0.18, fill=BLUE); box(s, 0, 7.32, 13.333, 0.18, fill=GREEN)
text(s, 0.9, 2.05, 11.5, 1.2, "OSINT Investigator", size=54, color=WHITE, bold=True)
text(s, 0.9, 3.15, 11.5, 0.8,
     "Automated external-OSINT investigation — from one seed to a cited intelligence graph",
     size=22, color=BLUE, bold=True)
text(s, 0.9, 4.3, 11.5, 1.6, [
    [("Give it a seed (domain, email, username, name, IP…) → it pivots through open sources, "
      "builds a confidence-tiered entity graph, and produces a cited CTI report.", 17, TEXT, False)],
    [("External OSINT only · Runs locally · Never fabricates a finding.", 16, GREEN, True)],
    [("Open source (MIT) · agent-vendor-agnostic (Claude Code or OpenAI Codex)", 14, MUTE, False)],
], space=8)
text(s, 0.9, 6.5, 11.5, 0.5, "Analyst briefing  ·  2026", size=13, color=MUTE)
notes(s, "A force-multiplier for an OSINT analyst. It does the tedious collection and first-pass link "
         "analysis at machine speed, and shows its work so you can trust or challenge every line. I'll "
         "finish with a real case where it mapped a ticket-fraud network from one URL.")

# ============================================================ SLIDE 2 — BLUF
s = slide(); header(s, "What it is", "One seed in → a cited intelligence graph out")
text(s, 0.55, 1.75, 6.0, 0.5, "THE PROBLEM", size=13, color=AMBER, bold=True)
text(s, 0.55, 2.1, 6.0, 1.2,
     "An analyst handed one selector faces hours of manual lookups across dozens of tools — then has "
     "to keep every link straight by hand.", size=16.5, color=TEXT)
text(s, 0.55, 3.35, 6.0, 0.5, "WHAT THIS DOES", size=13, color=GREEN, bold=True)
bullets(s, 0.55, 3.72, 6.2, 3.0, [
    "Runs the right tools automatically (ontology-driven, not hardcoded)",
    "Analyzes raw output and grades every finding by confidence",
    "Pivots on its own findings — loops until the leads run dry",
    "Red-teams itself before writing a report",
    "Ships a narrative report — every finding cited to tool output",
], size=15)
box(s, 7.0, 1.8, 5.75, 5.0, fill=PANEL, line=BLUE, line_w=1.5)
text(s, 7.3, 2.0, 5.15, 0.5, "THREE PROPERTIES THAT MAKE IT TRUSTWORTHY", size=13, color=BLUE, bold=True)
text(s, 7.3, 2.65, 5.15, 4.0, [
    [("1. No hallucination", 18, WHITE, True)],
    [("Every claim traces to a specific tool output or source URL.", 14.5, MUTE, False)],
    [(" ", 7, MUTE, False)],
    [("2. Full audit trail", 18, WHITE, True)],
    [("Every tool run logged raw; the graph is reproducible.", 14.5, MUTE, False)],
    [(" ", 7, MUTE, False)],
    [("3. Honest about limits", 18, WHITE, True)],
    [("Gaps are documented, never faked.", 14.5, MUTE, False)],
], space=5)
footer(s, "BLUF")
notes(s, "The headline for an analyst is the pivoting and confidence tiering. The headline for IT is "
         "the three properties on the right — local, auditable, external-only.")

# ============================================================ SLIDE 3 — OV-1 (reworked)
s = slide(); header(s, "Operational Concept (OV-1)", "Distinct agents, each a labeled SKILL — with an adversarial red team as its own stage")

def agentbox(x, y, w, h, title, skill, desc, accent):
    b = box(s, x, y, w, h, fill=PANEL, line=accent, line_w=1.75)
    text(s, x+0.12, y+0.08, w-0.24, 0.32, title, size=12.5, color=WHITE, bold=True)
    text(s, x+0.12, y+0.36, w-0.24, 0.28, "skill: " + skill, size=9.5, color=accent, bold=True)
    text(s, x+0.12, y+0.63, w-0.24, h-0.66, desc, size=9.5, color=MUTE)
    return b

# Band A: launch + brain
agentbox(0.55, 1.55, 1.75, 1.15, "LAUNCH", "investigate.md", "detect type · create case", BLUE)
arrow_r(s, 2.35, 2.0, 0.35)
brain = box(s, 2.8, 1.5, 2.55, 1.3, fill=PANEL, line=BLUE, line_w=2.25)
text(s, 2.9, 1.57, 2.4, 0.3, "SUPERVISOR", size=13.5, color=WHITE, bold=True)
text(s, 2.9, 1.88, 2.4, 0.26, "skill: supervisor.md", size=9.5, color=BLUE, bold=True)
text(s, 2.9, 2.14, 2.4, 0.6, "the analyst brain — routes via the ontology, analyzes, TIERS, "
     "commits the graph, pivots", size=9.5, color=MUTE)
text(s, 2.8, 2.85, 2.6, 0.3, "dispatches 3 collector agents ↓", size=10, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# Band B: THREE collector agents (distinct)
cy = 3.25
struct = agentbox(0.55, cy, 3.7, 1.35, "① STRUCTURED collector", "gatherer.md",
                  "58 typed tools — whois/rdap, dns, crt.sh, sherlock, shodan…", GREEN)
# paid CTI placeholder attached to the structured line
box(s, 0.7, cy+0.95, 3.4, 0.32, fill=PANEL2, line=AMBER, line_w=1.25, dash=True)
text(s, 0.78, cy+0.98, 3.3, 0.28, "▸ PAID CTI (roadmap): Shodan · DNSlytics · Flare",
     size=9.5, color=AMBER, bold=True)
agentbox(4.45, cy, 3.7, 1.35, "② WEB-SEARCH collector", "web_searcher.md",
         "real web searches + page fetches · snippet-as-evidence · cites everything", GREEN)
agentbox(8.35, cy, 4.35, 1.35, "③ ACTIVE collector", "active_collector.md",
         "touches the target's own site for tracker/analytics IDs (ownership corroborator) · passive-first OPSEC", GREEN)
text(s, 0.55, cy+1.42, 12.1, 0.3, "all three collect RAW · log everything · NEVER write the graph  "
     "(the raw / analysis split)", size=10.5, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

# Band C: analysis → gates → product
by = 5.35
box(s, 0.55, by, 2.75, 1.15, fill=PANEL, line=GREEN, line_w=1.75)
centered(box(s, 0.55, by, 2.75, 1.15, fill=None), [
    [("ANALYZE + TIER + COMMIT", 12, WHITE, True)],
    [("highly likely / probable / possible", 10.5, GREEN, True)],
    [("+ PIVOT each new entity (loop)", 10.5, TEXT, False)]])
arrow_r(s, 3.35, by+0.45, 0.4)
rt = box(s, 3.85, by, 3.05, 1.15, fill=PANEL, line=PURPLE, line_w=2.25)
centered(box(s, 3.85, by, 3.05, 1.15, fill=None), [
    [("RED TEAM  (separate agent)", 12, PURPLE, True)],
    [("skill: red_team.md", 10, PURPLE, True)],
    [("Mode 1: challenge the analysis", 10, TEXT, False)],
    [("Mode 2: ground the report", 10, TEXT, False)]])
arrow_r(s, 6.95, by+0.45, 0.4)
box(s, 7.45, by, 2.7, 1.15, fill=PANEL, line=BLUE, line_w=1.75)
centered(box(s, 7.45, by, 2.7, 1.15, fill=None), [
    [("REPORT WRITER", 12.5, WHITE, True)],
    [("skill: report-writer.md", 10, BLUE, True)],
    [("narrative CTI report", 10.5, TEXT, False)]])
arrow_r(s, 10.2, by+0.45, 0.35)
box(s, 10.6, by, 2.15, 1.15, fill=PANEL2, line=GREEN, line_w=1.75)
centered(box(s, 10.6, by, 2.15, 1.15, fill=None), [
    [("OUTPUTS", 12.5, GREEN, True)],
    [("graph · report", 10.5, TEXT, False)],
    [("bibliography · log", 10.5, TEXT, False)]])
text(s, 0.55, 6.62, 12.1, 0.4, "Separation of powers: collectors fetch · the supervisor alone analyzes & tiers · "
     "an INDEPENDENT red team tries to break the conclusions before they ship.",
     size=11.5, color=MUTE, align=PP_ALIGN.CENTER)
footer(s, "OV-1")
notes(s, "Walk left to right. Each box is a distinct agent running a named skill file. Three separate "
         "collectors — structured (typed tools; paid CTI like Shodan/DNSlytics/Flare plugs in here), "
         "web-search, and active. They only collect; the supervisor is the only brain that analyzes and "
         "tiers. The red team is its own agent — it challenges the analysis (Mode 1) and grounds the "
         "report (Mode 2) before anything ships.")

# ============================================================ SLIDE 4 — ONTOLOGY CATEGORIES
s = slide(); header(s, "The routing brain", "A knowledge base of selector CATEGORIES decides what runs — nothing hardcoded")
cats = [("Identity", "username · email · name · phone · handles", "28", True),
        ("Infrastructure", "domain · IP · URL · ASN · certificates", "11", True),
        ("Financial", "crypto wallets (BTC / ETH)", "5", True),
        ("Malware", "file hashes · IOCs", "3", True),
        ("Media", "images · EXIF / geolocation", "5", True),
        ("Geospatial", "coordinates · addresses", "3", True),
        ("Transport · Security · Social · Device", "plates · VIN · creds · devices", "10", False),
        ("+ uncatalogued roadmap types", "reachable via web-search fallback", "49", False)]
yy = 1.85
for i, (name, ex, cnt, live) in enumerate(cats):
    h = 0.6
    box(s, 0.55, yy, 8.7, h, fill=PANEL if i % 2 == 0 else PANEL2)
    dot = GREEN if live else MUTE
    box(s, 0.7, yy+0.19, 0.22, 0.22, fill=dot, shape=MSO_SHAPE.OVAL)
    text(s, 1.05, yy+0.06, 5.0, 0.5, name, size=14.5, color=WHITE, bold=True)
    text(s, 1.05, yy+0.32, 6.3, 0.28, ex, size=10.5, color=MUTE)
    text(s, 7.5, yy+0.06, 1.6, 0.5, cnt + " types", size=13, color=(GREEN if live else MUTE),
         bold=True, align=PP_ALIGN.RIGHT)
    yy += h + 0.04
box(s, 9.55, 1.85, 3.2, 5.0, fill=PANEL, line=BLUE, line_w=1.5)
text(s, 9.78, 2.05, 2.75, 4.7, [
    [("92", 34, BLUE, True)],
    [("selector types the ontology can reason about", 12.5, TEXT, False)],
    [(" ", 8, MUTE, False)],
    [("58", 34, GREEN, True)],
    [("runnable tools wired in", 12.5, TEXT, False)],
    [(" ", 8, MUTE, False)],
    [("● live coverage", 11.5, GREEN, True)],
    [("● roadmap", 11.5, MUTE, True)],
    [(" ", 8, MUTE, False)],
    [("Routing: plan_collection(). Pivoting: each type declares what it YIELDS. A test enforces "
      "the ontology can't lie about what runs.", 11.5, MUTE, False)],
], space=3)
footer(s, "Ontology → categories")
notes(s, "The system reasons in categories. Identity and Infrastructure are where the live tool coverage "
         "is deepest — exactly what most investigations need. The rest is a roadmap, and a bare selector "
         "always has the web-search line as a fallback. A regression test keeps the self-description honest.")

# ============================================================ SLIDE 5 — TOOLS SAMPLE (real)
s = slide(); header(s, "The tools", "A sample of the 58 tools running now — real names, grouped by selector")
groups = [
    ("DOMAIN (16)", GREEN, "whois_lookup · rdap · dns_lookup · dnsrecon · crtsh · certspotter · tls_cert · "
     "http_title · urlscan · wayback · theharvester · cloud_buckets · web_tech_fingerprint"),
    ("IP (11)", BLUE, "ip_geolocation · ipinfo · shodan_internetdb · reverse_dns · ripestat_network · "
     "bgpview_ip · greynoise_community · reverse_ip · robtex_ip · urlscan"),
    ("USERNAME (8)", PURPLE, "sherlock · maigret · naminter · linkook · socialscan · github_user · "
     "reddit_about · google_dork_generator"),
    ("EMAIL (8)", PINK, "holehe · user_scanner · hudsonrock_email · xposedornot_email · disify · "
     "socialscan · gravatar_check · pgp_keyserver"),
    ("NAME / COMPANY (11)", AMBER, "wikipedia_search · wikidata_search · gleif_lei · sec_edgar_fts · "
     "aleph_occrp · courtlistener_search · hibp_name_search · name_to_username"),
    ("PHONE · CRYPTO · IMAGE · TRACKER", TEXT, "phonenumbers · ignorant · phoneinfoga · blockstream_btc · "
     "blockchain_btc · exiftool · web_tech_fingerprint · tracker_reverse"),
]
yy = 1.8
for name, col, tools in groups:
    box(s, 0.55, yy, 12.2, 0.8, fill=PANEL if groups.index((name, col, tools)) % 2 == 0 else PANEL2)
    text(s, 0.72, yy+0.09, 3.1, 0.6, name, size=13, color=col, bold=True)
    text(s, 3.7, yy+0.08, 8.9, 0.65, tools, size=11.5, color=TEXT)
    yy += 0.84
text(s, 0.55, 6.75, 12.2, 0.4, "Each tool is a small declarative spec (endpoint + parser). Optional CLIs "
     "and keyed APIs degrade gracefully — the system runs with zero keys.", size=11, color=MUTE)
footer(s, "Tools in use")
notes(s, "These are real tool names in the registry today, grouped by what they run against. Best-in-class "
         "on the core selectors — domain, IP, username, email, name. Each is a readable spec, so IT can "
         "see exactly what every tool touches.")

# ============================================================ SLIDE 6 — SOFTWARE LAYDOWN
s = slide(); header(s, "Software laydown", "Small, auditable, standard — pure-Python engine")
text(s, 0.55, 1.8, 6.0, 0.4, "CORE DEPENDENCIES (7, pinned in requirements.txt)", size=13, color=GREEN, bold=True)
deps = [("requests", "HTTP client for all API tools"), ("networkx", "the investigation graph"),
        ("python-whois", "WHOIS lookups"), ("dnspython", "DNS records"),
        ("beautifulsoup4", "HTML parsing"), ("mmh3", "favicon hashing (ownership corroborator)"),
        ("phonenumbers", "phone parsing/validation")]
yy = 2.25
for i, (pk, ds) in enumerate(deps):
    box(s, 0.55, yy, 6.0, 0.5, fill=PANEL if i % 2 == 0 else PANEL2)
    text(s, 0.7, yy+0.08, 2.3, 0.4, pk, size=13.5, color=BLUE, bold=True)
    text(s, 3.0, yy+0.09, 3.4, 0.4, ds, size=12, color=MUTE)
    yy += 0.54
box(s, 6.95, 1.8, 5.8, 5.0, fill=PANEL, line=GREEN, line_w=1.5)
text(s, 7.2, 1.95, 5.3, 0.4, "WHAT WE BUILT  (~48 Python modules under src/)", size=13, color=GREEN, bold=True)
text(s, 7.2, 2.45, 5.35, 4.3, [
    [("core/", 14, WHITE, True), ("  selector detection + investigation state", 13, MUTE, False)],
    [("ontology/", 14, WHITE, True), ("  the routing brain + honesty annotator", 13, MUTE, False)],
    [("tools/", 14, WHITE, True), ("  58 tools via HttpTool/CliTool specs + registry", 13, MUTE, False)],
    [("graph/", 14, WHITE, True), ("  NetworkX graph + vis.js visualizer", 13, MUTE, False)],
    [("report/", 14, WHITE, True), ("  narrative report (MD + HTML + Mermaid)", 13, MUTE, False)],
    [("logger/", 14, WHITE, True), ("  the raw audit trail", 13, MUTE, False)],
    [("scripts/health_check.py", 14, GREEN, True), ("  the safety gate (registry + floor + 3 suites)", 13, MUTE, False)],
    [(" ", 8, MUTE, False)],
    [("One command — ./bootstrap.sh — builds and self-verifies. Full install list in the appendix.",
      12.5, AMBER, False)],
], space=6)
footer(s, "Software laydown")
notes(s, "Seven mainstream open-source packages, everything pinned, one command to build and self-verify. "
         "No black box — the whole dependency list is on one slide, and the exact install list is in the appendix.")

# ============================================================ SLIDE 7 — USE CASE 1 (overview image)
s = slide(); header(s, "Use case · ticket-fraud network", "One scam URL → a mapped network (auto-generated)")
box(s, 0.55, 1.5, 4.3, 5.15, fill=PANEL, line=BLUE, line_w=1.5)
text(s, 0.75, 1.65, 3.95, 5.0, [
    [("SEED", 12, MUTE, True)],
    [("colosseumdiroma-tickets.com", 14, BLUE, True)],
    [("a fake Colosseum ticket-resale site", 12, TEXT, False)],
    [(" ", 8, MUTE, False)],
    [("WHAT THE PIVOTS BUILT", 12, GREEN, True)],
    [("• Registration — WHOIS/RDAP → registrant org 'The Walker Tours LLC'", 12, TEXT, False)],
    [("• Hosting estate — domain → AWS IPs → co-hosted lookalikes → loop", 12, TEXT, False)],
    [("• Content IDs — the active line pulled 15 tracker/analytics IDs", 12, TEXT, False)],
    [(" ", 6, MUTE, False)],
    [("TIERED HONESTLY", 12, GREEN, True)],
    [("96 highly likely", 14, GREEN, True), (" · ", 12, MUTE, False), ("14 probable", 12.5, AMBER, True),
     (" · ", 12, MUTE, False), ("9 possible", 12.5, MUTE, True)],
    [("Weak leads kept as pivots, not dropped. NOT over-sold as a 75-site scam estate — "
      "AWS nameservers & the legit origin domains are flagged.", 11, MUTE, False)],
], space=5)
pic(s, os.path.join(IMG, "graph_overview.png"), 5.05, 1.5, w=7.75)
text(s, 5.05, 6.62, 7.75, 0.35, "Real output: 119 entities / 170 relationships, colored by type "
     "(blue=domain, green=IP hub, amber=company, purple=tracker ID).", size=10.5, color=MUTE, align=PP_ALIGN.CENTER)
footer(s, "Use case 1/3")
notes(s, "Start at one URL. The system expanded the hosting neighborhood, fingerprinted the sites, and "
         "clustered them — automatically. This is the ACTUAL graph it produced: 119 entities from one seed. "
         "Note the two big green nodes — shared hosting IPs — and the amber company nodes where attribution lands.")

# ============================================================ SLIDE 8 — USE CASE 2 (attribution + tracker images)
s = slide(); header(s, "Use case · what it proves", "Who runs it — and the independent evidence that ties them together")
pic(s, os.path.join(IMG, "graph_attribution.png"), 0.4, 1.55, w=7.5)
text(s, 0.4, 5.75, 7.5, 0.35, "‘Who runs it’: real operator entities the pivots recovered — Walker Tours "
     "LLC/Corp., Pancho Tours, Grupo Feel The City, shared EIN 37-2091569, CIFs, names, emails, phones.",
     size=10, color=MUTE)
pic(s, os.path.join(IMG, "graph_tracker.png"), 8.05, 1.75, w=4.85)
text(s, 8.05, 5.35, 4.85, 0.8, [
    [("THE CORROBORATOR MECHANISM", 12, PURPLE, True)],
    [("One shared Google Ads ID (AW-778701917) links 8 scam domains. A shared strong tracker ID is the "
      "independent evidence that upgrades ‘co-hosted’ to ‘same operator’.", 11, TEXT, False)],
], space=3)
footer(s, "Use case 2/3")
notes(s, "Left: the attribution cluster — real companies, a shared federal EIN, Spanish CIFs, names, "
         "contact emails and phones the system recovered. Right: HOW it knows — a single shared analytics "
         "ID tying eight domains together. That shared ID is the independent corroborator our doctrine "
         "requires before claiming common ownership.")

# ============================================================ SLIDE 9 — USE CASE 3 (trust story)
s = slide(); header(s, "Use case · the trust story", "The system pressure-tests itself and refuses to over-claim")
steps = [
    ("1", "An early pass read the estate as TWO operators.", TEXT),
    ("2", "The RED-TEAM gate challenged the merge: ‘shared hosting ≠ shared ownership — show an independent corroborator.’", TEXT),
    ("3", "A deeper pivot found the GOLD STANDARD: the seed itself ran a strong UA-131208121-1 analytics property in its 2024 Wayback snapshot — the same property as sevillafreetour.com.", GREEN),
    ("4", "On SIX independent corroborators, red-team round 3 REVISED the conclusion to ONE operator group: Grupo Feel The City S.L. (Seville), fronted by two US shells sharing federal EIN 37-2091569.", TEXT),
    ("5", "It distinguished the LEGITIMATE 2010 origin tour business from the 2024–26 scam build-out — and flagged the legit brand must NOT be labelled a scam.", AMBER),
]
yy = 1.85
for num, txt, col in steps:
    c = box(s, 0.55, yy, 0.5, 0.5, fill=PANEL, line=BLUE, line_w=1.5, shape=MSO_SHAPE.OVAL)
    centered(c, [[(num, 15, BLUE, True)]])
    text(s, 1.3, yy-0.02, 11.4, 0.9, txt, size=14, color=col, anchor=MSO_ANCHOR.MIDDLE)
    yy += 0.92
box(s, 0.55, 6.5, 12.2, 0.55, fill=PANEL, line=GREEN, line_w=1.5)
text(s, 0.75, 6.59, 11.9, 0.4, "Why it matters: it argued with itself on the record, changed its mind on "
     "evidence, and drew a clean line around the innocent business. Every step is in the audit log.",
     size=13, color=WHITE, bold=True)
footer(s, "Use case 3/3")
notes(s, "If a tool only ever confirms its first guess, you can't trust it. This one argued with itself on "
         "the record, upgraded its conclusion when the evidence justified it, and protected the innocent "
         "business. That discipline is what makes machine-assisted attribution safe to brief.")

# ============================================================ SLIDE 10 — IT 1 (egress)
s = slide(); header(s, "Why IT should approve · 1", "Where the data goes")
text(s, 0.55, 1.75, 12, 0.5, "CONCERN:  data egress / data handling", size=15, color=AMBER, bold=True)
bullets(s, 0.55, 2.35, 12.2, 4.2, [
    [("Runs entirely on the local machine.", 16, WHITE, True), ("  No server, no cloud service, no account required to operate.", 15, TEXT, False)],
    [("External OSINT only — by design and scope.", 16, WHITE, True), ("  Never ingests internal/corporate data; internal-data analysis is explicitly out of scope.", 15, TEXT, False)],
    [("What leaves the machine is only the selector,", 16, WHITE, True), ("  sent to PUBLIC sources (WHOIS/DNS/CT logs, search engines, public APIs) — the same requests an analyst makes in a browser.", 15, TEXT, False)],
    [("Results stay on disk.", 16, WHITE, True), ("  investigations/ and the .env key file are git-ignored — case data and keys are never committed or pushed.", 15, TEXT, False)],
    [("Passive-first OPSEC.", 16, WHITE, True), ("  The active line prefers archived (Wayback) copies before any live request; generic user-agent; no crawling; proxy seam for attribution control.", 15, TEXT, False)],
], size=15, gap=11)
box(s, 0.55, 6.5, 12.2, 0.55, fill=PANEL, line=GREEN, line_w=1.5)
text(s, 0.75, 6.59, 11.9, 0.4, "Net: the tool's network behavior is a subset of normal analyst browsing, "
     "from the local machine, against public sources only.", size=13, color=WHITE, bold=True)
footer(s, "IT approval · data egress")
notes(s, "Nothing internal goes in. What goes out is the same public lookups an analyst runs manually. "
         "Results and keys never leave the box and are never committed to git.")

# ============================================================ SLIDE 11 — IT 2 (provenance)
s = slide(); header(s, "Why IT should approve · 2", "Dependency & tool provenance (supply chain)")
bullets(s, 0.55, 1.8, 12.2, 3.3, [
    [("Seven pinned, mainstream open-source packages.", 15.5, WHITE, True), ("  The entire core dependency list fits on one slide (requirements.txt). No obscure runtime.", 14.5, TEXT, False)],
    [("Declarative tool registry.", 15.5, WHITE, True), ("  Every one of the 58 tools is a small, readable spec (endpoint + parser) — auditable at a glance; no hidden collection.", 14.5, TEXT, False)],
    [("No bespoke scraping allowed.", 15.5, WHITE, True), ("  Collection goes only through registered, reviewed tools; the doctrine forbids ad-hoc scripts.", 14.5, TEXT, False)],
    [("Optional extras are opt-in and degrade gracefully.", 15.5, WHITE, True), ("  Fully functional with zero keys and zero extra installs.", 14.5, TEXT, False)],
    [("MIT-licensed", 15.5, WHITE, True), (" — open for review and download.", 14.5, TEXT, False)],
], size=14.5, gap=9)
box(s, 0.55, 5.25, 12.2, 1.4, fill=PANEL, line=BLUE, line_w=1.5)
text(s, 0.8, 5.4, 11.7, 0.4, "BUILT-IN SAFETY GOVERNANCE (bonus)", size=13, color=BLUE, bold=True)
text(s, 0.8, 5.8, 11.7, 0.8, "A health gate (health_check.py) + 3 regression suites gate every change · "
     "a capability-lock defines what must never regress · a separate System Manager role maintains it conservatively.",
     size=14, color=TEXT)
footer(s, "IT approval · provenance")
notes(s, "Everything a supply-chain reviewer wants: short pinned dependency list, every tool is a readable "
         "spec, no hidden network calls, opt-in extras, permissive license, and an automated gate that "
         "proves the system still works after any change.")

# ============================================================ SLIDE 12 — CLOSE
s = slide(); header(s, "Close", "Ready, open, and built to be trusted")
text(s, 0.55, 1.85, 12, 0.4, "WHAT THIS GIVES THE TEAM", size=13, color=GREEN, bold=True)
bullets(s, 0.55, 2.25, 12.2, 1.8, [
    "An analyst force-multiplier: seed → cited, confidence-tiered graph + report, automatically, with a full audit trail",
    "Proven on a real ticket-fraud network — mapped, attributed, and self-corrected without over-claiming",
    "Safe by construction: local · external-only · auditable · honest about limits",
], size=15.5, gap=8)
box(s, 0.55, 4.25, 12.2, 1.5, fill=PANEL, line=BLUE, line_w=1.5)
text(s, 0.8, 4.4, 11.7, 1.3, [
    [("AVAILABLE NOW", 13, BLUE, True)],
    [("Open-source (MIT):  github.com/cis379/osint-investigator-v3  (public)", 15, TEXT, False)],
    [("One command to stand up (./bootstrap.sh), health-gated, Mac / Linux / Windows · drive it with Claude Code or Codex", 14.5, TEXT, False)],
], space=5)
text(s, 0.55, 6.0, 12, 0.7, [
    [("THE ASK:  ", 18, AMBER, True),
     ("approval to run it locally, external-OSINT only, as an analyst tool.", 18, WHITE, True)]])
footer(s, "Close")
notes(s, "It's ready, it's open, and it's built to be trusted. I'm asking for the green light to use it as "
         "a local analyst tool — and I'm happy to walk IT through the code and the audit trail.")

# ============================================================ SLIDE 13 — APPENDIX: INSTALL
s = slide(); header(s, "Appendix", "Exactly what must be installed (full software footprint)")
# Required
box(s, 0.55, 1.7, 6.0, 2.55, fill=PANEL, line=GREEN, line_w=1.5)
text(s, 0.75, 1.82, 5.6, 0.35, "REQUIRED", size=13, color=GREEN, bold=True)
text(s, 0.75, 2.2, 5.6, 2.0, [
    [("Python 3.10+", 14, WHITE, True)],
    [("7 pip packages (requirements.txt):", 13, TEXT, True)],
    [("requests · networkx · python-whois · dnspython · beautifulsoup4 · mmh3 · phonenumbers", 12, MUTE, False)],
    [(" ", 6, MUTE, False)],
    [("→ one command: ", 13, TEXT, False), ("./bootstrap.sh", 13, GREEN, True),
     (" does all of the below + a health check", 12.5, MUTE, False)],
], space=5)
# Optional CLIs
box(s, 6.75, 1.7, 6.0, 2.55, fill=PANEL, line=BLUE, line_w=1.5)
text(s, 6.95, 1.82, 5.6, 0.35, "OPTIONAL — external OSINT CLIs (degrade gracefully)", size=12.5, color=BLUE, bold=True)
text(s, 6.95, 2.25, 5.65, 2.0, [
    [("pipx install:", 13, WHITE, True)],
    [("sherlock-project · maigret · holehe · theHarvester · socialscan · dnsrecon · socid-extractor · user-scanner · linkook", 11.5, MUTE, False)],
    [("pip (into venv):", 13, WHITE, True), ("  naminter · ignorant", 12, MUTE, False)],
    [("system pkg:", 13, WHITE, True), ("  exiftool  (brew / apt / winget)", 12, MUTE, False)],
], space=5)
# Optional keyed
box(s, 0.55, 4.4, 12.2, 1.55, fill=PANEL2, line=AMBER, line_w=1.5)
text(s, 0.75, 4.52, 11.8, 0.35, "OPTIONAL — keyed APIs (.env, OFF by default; the system runs fully without any)", size=12.5, color=AMBER, bold=True)
text(s, 0.75, 4.92, 11.8, 1.0, [
    [("Paid/keyed CTI that plugs into the STRUCTURED collector:  ", 13, TEXT, True),
     ("Shodan · DNSlytics · Flare · VirusTotal · AbuseIPDB · GreyNoise", 13, WHITE, False)],
    [("Each degrades to a clear ‘needs KEY’ skip when absent — no key is required to operate.", 12, MUTE, False)],
], space=5)
text(s, 0.55, 6.15, 12.2, 0.7, "Footprint summary: one Python venv + up to ~11 well-known OSINT CLIs + "
     "one binary (exiftool). No services, no daemons, no admin rights required.", size=12.5, color=TEXT)
footer(s, "Appendix · install")
notes(s, "This is the complete software footprint for IT. Required is just Python plus seven packages. "
         "Everything else is optional and degrades gracefully. bootstrap.sh installs it all and runs the "
         "health check. Paid CTI keys (Shodan/DNSlytics/Flare) are opt-in and plug into the structured collector.")

out = os.path.join(HERE, "osint-brief.pptx")
prs.save(out)
print("WROTE", out, "-", len(prs.slides._sldIdLst), "slides")
